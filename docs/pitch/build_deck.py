"""Build the ElimuTayari lightning-pitch deck (3 minutes, 6 slides).

Written for hackathon judges: one claim per slide, type large enough to read from
the back of a room, and a demo hand-off slide that says exactly what they are
about to watch.

    python3 docs/pitch/build_deck.py

Placeholders to fill before presenting are collected in PRESENTER at the top.
"""

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUT = "docs/pitch/elimutayari-pitch.pptx"

PRESENTER = {
    "names": "Nirel · Alphonce · David · Millicent · Rhodah",
    "event": "AI Mashinani",
    "ussd_code": "*384*XXXX#",
}

# Ink on paper, with a single green accent. Deliberately calm: the demo is the
# colour in this pitch, not the slides.
INK = RGBColor(0x14, 0x1A, 0x16)
BODY = RGBColor(0x44, 0x4E, 0x47)
MUTED = RGBColor(0x7E, 0x8A, 0x82)
PAPER = RGBColor(0xFA, 0xFA, 0xF7)
GREEN = RGBColor(0x0E, 0x7A, 0x4E)
GREEN_SOFT = RGBColor(0xE4, 0xF0, 0xE9)
RULE = RGBColor(0xDC, 0xE2, 0xDD)

HEAD = "Verdana"   # ships on Windows and macOS, so the deck opens as designed
TEXT = "Verdana"

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.85)


def deck():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    return prs


def blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PAPER
    bg.line.fill.background()
    bg.shadow.inherit = False
    return slide


def textbox(slide, x, y, w, h, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    frame.paragraphs[0].alignment = align
    return frame


def write(frame, text, size, color=INK, bold=False, font=TEXT, space_after=0,
          line=None, align=None, first=False):
    para = frame.paragraphs[0] if first else frame.add_paragraph()
    para.text = text
    para.space_after = Pt(space_after)
    if line:
        para.line_spacing = line
    if align:
        para.alignment = align
    for run in para.runs:
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font
        run.font.color.rgb = color
    return para


def rule(slide, x, y, w, color=RULE, thickness=Pt(1.25)):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, thickness)
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    line.shadow.inherit = False
    return line


def card(slide, x, y, w, h, fill=None, edge=RULE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.06
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill or PAPER
    shape.line.color.rgb = edge
    shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape


def chrome(slide, label, number):
    """Section label top-left, slide number bottom-right, hairline under the head."""
    frame = textbox(slide, MARGIN, Inches(0.5), Inches(8), Inches(0.3))
    write(frame, label.upper(), 11, GREEN, bold=True, font=HEAD, first=True)
    frame = textbox(slide, W - MARGIN - Inches(1), H - Inches(0.72), Inches(1),
                    Inches(0.3), align=PP_ALIGN.RIGHT)
    write(frame, str(number), 11, MUTED, font=HEAD, first=True)


def heading(slide, text, size=40):
    frame = textbox(slide, MARGIN, Inches(0.95), W - 2 * MARGIN, Inches(1.1))
    write(frame, text, size, INK, bold=True, font=HEAD, line=0.95, first=True)
    rule(slide, MARGIN, Inches(2.0), Inches(1.6), GREEN, Pt(3))


# ---------------------------------------------------------------- slide 1


def title_slide(prs):
    slide = blank(prs)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.28), H)
    band.fill.solid()
    band.fill.fore_color.rgb = GREEN
    band.line.fill.background()
    band.shadow.inherit = False

    frame = textbox(slide, MARGIN, Inches(2.05), Inches(10.6), Inches(1.2))
    write(frame, "ElimuTayari", 66, INK, bold=True, font=HEAD, first=True)

    frame = textbox(slide, MARGIN, Inches(3.25), Inches(10.2), Inches(1.5))
    write(frame, "A CBC teaching companion for every Grade 10 teacher in Kenya —"
                 " on the phone they already own.", 25, BODY, font=TEXT, line=1.25,
          first=True)

    rule(slide, MARGIN, Inches(4.85), Inches(2.2), GREEN, Pt(3))

    frame = textbox(slide, MARGIN, Inches(5.25), Inches(10), Inches(1.1))
    write(frame, "No smartphone. No data bundle. No app to install.", 18, INK,
          bold=True, font=TEXT, space_after=8, first=True)
    write(frame, PRESENTER["names"], 14, MUTED, font=TEXT, space_after=4)
    write(frame, PRESENTER["event"], 14, GREEN, bold=True, font=TEXT)
    return slide


# ---------------------------------------------------------------- slide 2


def problem_slide(prs):
    slide = blank(prs)
    chrome(slide, "The problem", 2)
    heading(slide, "Senior School CBC landed. The support did not.")

    frame = textbox(slide, MARGIN, Inches(2.4), Inches(6.5), Inches(3.4))
    write(frame, "Grade 10 began under a curriculum most teachers had never taught.",
          22, INK, bold=True, font=TEXT, space_after=16, line=1.2, first=True)
    write(frame, "Each learning area is a new KICD design: English alone is 45 "
                 "sub-strands across 9 themes and 180 lessons, each with its own "
                 "outcomes, learning experiences and assessment.", 16, BODY,
          font=TEXT, space_after=14, line=1.35)
    write(frame, "The teaching aids built for it assume a smartphone, a data bundle "
                 "and an app store — which is not what a teacher in Turkana or "
                 "Kitui is holding.", 16, BODY, font=TEXT, line=1.35)

    # The gap, stated as the two facts that matter, side by side.
    x = Inches(7.9)
    for i, (stat, label) in enumerate([
        ("100%", "of Kenyan teachers can be reached on a basic GSM phone — USSD "
                 "and SMS need no data and no app"),
        ("0", "curriculum-faithful teaching packs delivered that way today"),
    ]):
        y = Inches(2.4) + i * Inches(1.75)
        card(slide, x, y, Inches(4.55), Inches(1.5), GREEN_SOFT, GREEN_SOFT)
        frame = textbox(slide, x + Inches(0.35), y + Inches(0.22), Inches(1.35),
                        Inches(1.1), anchor=MSO_ANCHOR.MIDDLE)
        write(frame, stat, 36, GREEN, bold=True, font=HEAD, first=True)
        frame = textbox(slide, x + Inches(1.75), y + Inches(0.22), Inches(2.6),
                        Inches(1.1), anchor=MSO_ANCHOR.MIDDLE)
        write(frame, label, 13, INK, font=TEXT, line=1.25, first=True)

    frame = textbox(slide, x, Inches(5.95), Inches(4.55), Inches(0.9))
    write(frame, "So the phone is not the constraint. The delivery is.", 15, GREEN,
          bold=True, font=TEXT, line=1.3, first=True)
    return slide


# ---------------------------------------------------------------- slide 3


def how_slide(prs):
    slide = blank(prs)
    chrome(slide, "How it works", 3)
    heading(slide, "Dial a code. Teach the lesson. Reply with what confused them.")

    steps = [
        ("1", "DIAL", f"Teacher dials {PRESENTER['ussd_code']} and browses\n"
                      "learning area → strand → sub-strand"),
        ("2", "RECEIVE", "Pack arrives by SMS in whole\n160-character parts"),
        ("3", "REPLY", "After class: Q M-ALG-02 Why do we\nswap the sign?"),
        ("4", "COMPOUND", "Questions drive test generation and\ntopic-difficulty "
                          "analytics"),
    ]

    x = MARGIN
    width = Inches(2.62)
    gap = Inches(0.28)
    for i, (num, label, body) in enumerate(steps):
        left = x + i * (width + gap)
        shape = card(slide, left, Inches(2.5), width, Inches(2.35))
        shape.line.color.rgb = GREEN if i in (0, 2) else RULE

        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.28),
                                       Inches(2.75), Inches(0.46), Inches(0.46))
        badge.fill.solid()
        badge.fill.fore_color.rgb = GREEN
        badge.line.fill.background()
        badge.shadow.inherit = False
        frame = badge.text_frame
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        write(frame, num, 15, PAPER, bold=True, font=HEAD, align=PP_ALIGN.CENTER,
              first=True)

        frame = textbox(slide, left + Inches(0.28), Inches(3.4), width - Inches(0.56),
                        Inches(0.3))
        write(frame, label, 12, GREEN, bold=True, font=HEAD, first=True)
        frame = textbox(slide, left + Inches(0.28), Inches(3.78),
                        width - Inches(0.56), Inches(1.0))
        for j, line in enumerate(body.split("\n")):
            write(frame, line, 13, BODY, font=TEXT, line=1.3, first=(j == 0))

        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, left + width + Inches(0.04), Inches(3.55),
                Inches(0.2), Inches(0.2))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MUTED
            arrow.line.fill.background()
            arrow.shadow.inherit = False

    loop = card(slide, MARGIN, Inches(5.2), W - 2 * MARGIN, Inches(1.05), GREEN_SOFT,
                GREEN_SOFT)
    frame = loop.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = frame.margin_right = Inches(0.35)
    write(frame, "The loop closes: every question a teacher sends back makes the "
                 "next teacher's pack, and the next test, better.", 17, INK,
          bold=True, font=TEXT, align=PP_ALIGN.CENTER, first=True)
    return slide


# ---------------------------------------------------------------- slide 4


def demo_slide(prs):
    slide = blank(prs)
    chrome(slide, "Live demo", 4)
    heading(slide, "Watch this happen on a real phone.")

    frame = textbox(slide, MARGIN, Inches(2.35), Inches(6.3), Inches(3.6))
    write(frame, "What you are about to see", 15, GREEN, bold=True, font=HEAD,
          space_after=14, first=True)
    for step in [
        f"Dial {PRESENTER['ussd_code']} — no data, no app.",
        "Browse Core Mathematics → Algebra → Indices and Logarithms.",
        "The teaching pack arrives by SMS, split into whole 160-character parts.",
        "Reply Q M-ALG-02 with a student's question.",
        "The question lands against that sub-strand, ready for test generation.",
    ]:
        write(frame, "→  " + step, 16, BODY, font=TEXT, space_after=11, line=1.3)

    # The USSD screen, as the state machine actually renders it.
    phone = card(slide, Inches(7.9), Inches(2.35), Inches(4.55), Inches(3.15),
                 RGBColor(0x14, 0x1A, 0x16), RGBColor(0x14, 0x1A, 0x16))
    frame = phone.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.32)
    frame.margin_top = Inches(0.3)
    screen = [
        "Welcome to ElimuTayari",
        "1. Core Mathematics",
        "Or enter a sub-strand code e.g. M-ALG-02",
    ]
    for i, line in enumerate(screen):
        write(frame, line, 16, RGBColor(0x8B, 0xE8, 0xB4) if i == 0 else PAPER,
              bold=(i == 0), font="Courier New", space_after=9, line=1.25,
              first=(i == 0))
    write(frame, "—", 14, MUTED, font="Courier New", space_after=9)
    write(frame, "Q M-ALG-02 Why do we swap the sign?", 15,
          RGBColor(0x8B, 0xE8, 0xB4), font="Courier New", line=1.25)

    frame = textbox(slide, Inches(7.9), Inches(5.65), Inches(4.55), Inches(0.7))
    write(frame, "Every screen fits 160 characters — enforced by test, not by "
                 "hope.", 13, MUTED, font=TEXT, line=1.3, first=True)
    return slide


# ---------------------------------------------------------------- slide 5


def architecture_slide(prs):
    """The core engineering decision: a pre-generated wiki, not a live chatbot."""
    slide = blank(prs)
    chrome(slide, "The architecture decision", 5)
    heading(slide, "We chose a wiki, not a chatbot.")

    frame = textbox(slide, MARGIN, Inches(2.3), Inches(11.6), Inches(1.0))
    write(frame, "Curriculum is static. A sub-strand does not change per teacher or "
                 "per question — so generating it once and reviewing it once beats "
                 "regenerating it on every request.", 18, INK, font=TEXT, line=1.3,
          first=True)

    items = [
        ("Zero token cost per lesson",
         "A teaching pack is a database read, not an inference. The next teacher "
         "costs one SMS — so this scales without the model bill scaling with it."),
        ("Deterministic, so reviewable",
         "The same code returns the same page every time. That is what makes "
         "sign-off possible — you cannot review an answer regenerated on every ask."),
        ("Hallucination cannot reach a class",
         "Content comes from the official KICD design, marked draft-human-review "
         "until a teacher approves it. The live model only writes tests and "
         "clusters questions."),
    ]
    width = Inches(3.65)
    gap = Inches(0.3)
    for i, (title, body) in enumerate(items):
        left = MARGIN + i * (width + gap)
        card(slide, left, Inches(3.35), width, Inches(2.85))
        rule(slide, left + Inches(0.3), Inches(3.65), Inches(0.75), GREEN, Pt(3))
        frame = textbox(slide, left + Inches(0.3), Inches(3.88), width - Inches(0.6),
                        Inches(0.62))
        write(frame, title, 16, INK, bold=True, font=HEAD, line=1.15, first=True)
        frame = textbox(slide, left + Inches(0.3), Inches(4.62), width - Inches(0.6),
                        Inches(1.55))
        write(frame, body, 13, BODY, font=TEXT, line=1.35, first=True)

    frame = textbox(slide, MARGIN, Inches(6.35), Inches(11.6), Inches(0.4))
    write(frame, "FastAPI · Alembic migrations that run on SQLite and Postgres alike "
                 "· injectable Claude seam · 72 automated tests", 12, MUTED,
          font=TEXT, first=True)
    return slide


# ---------------------------------------------------------------- slide 6


def channels_slide(prs):
    """Online and offline are different jobs for different people."""
    slide = blank(prs)
    chrome(slide, "Online and offline by design", 6)
    heading(slide, "The board works online. The teacher does not have to.")

    bands = [
        ("ONLINE", "The education board owns the content", GREEN, Inches(2.3), [
            "Board reviews the official KICD design",
            "Wiki updated — Markdown, versioned, signed off",
            "Seeded into the database on deploy",
        ]),
        ("OFFLINE", "The teacher needs no connectivity", INK, Inches(4.65), [
            "Dials USSD from any GSM phone",
            "Teaching pack arrives by SMS",
            "Replies with student questions",
        ]),
    ]

    for label, caption, color, top, steps in bands:
        frame = textbox(slide, MARGIN, top, Inches(3.4), Inches(0.3))
        write(frame, label, 13, color, bold=True, font=HEAD, first=True)
        frame = textbox(slide, MARGIN + Inches(1.35), top, Inches(6.5), Inches(0.3))
        write(frame, caption, 13, MUTED, font=TEXT, first=True)

        width = Inches(3.65)
        gap = Inches(0.3)
        for i, step in enumerate(steps):
            left = MARGIN + i * (width + gap)
            fill = GREEN_SOFT if label == "ONLINE" else PAPER
            shape = card(slide, left, top + Inches(0.42), width, Inches(1.25), fill,
                         GREEN_SOFT if label == "ONLINE" else RULE)
            frame = shape.text_frame
            frame.word_wrap = True
            frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            frame.margin_left = frame.margin_right = Inches(0.28)
            write(frame, step, 15, INK, font=TEXT, line=1.3, first=True)

            if i < len(steps) - 1:
                arrow = slide.shapes.add_shape(
                    MSO_SHAPE.RIGHT_ARROW, left + width + Inches(0.05),
                    top + Inches(0.96), Inches(0.2), Inches(0.18))
                arrow.fill.solid()
                arrow.fill.fore_color.rgb = MUTED
                arrow.line.fill.background()
                arrow.shadow.inherit = False

    # The seam between the two: where the internet requirement stops.
    rule(slide, MARGIN, Inches(4.28), W - 2 * MARGIN, RULE, Pt(1.25))
    seam = card(slide, Inches(3.9), Inches(4.05), Inches(5.55), Inches(0.46), PAPER,
                GREEN)
    frame = seam.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    write(frame, "Below this line, nothing needs the internet", 13, GREEN, bold=True,
          font=TEXT, align=PP_ALIGN.CENTER, first=True)

    frame = textbox(slide, MARGIN, Inches(6.35), Inches(11.6), Inches(0.5))
    write(frame, "Content updates are an institutional job, done once, online. "
                 "Teaching is a daily job, done anywhere, offline.", 15, INK,
          bold=True, font=TEXT, first=True)
    return slide


# ---------------------------------------------------------------- slide 6


def ask_slide(prs):
    slide = blank(prs)
    chrome(slide, "Where this goes", 7)
    heading(slide, "One subject works. The pipeline does the rest.")

    columns = [
        ("BUILT", GREEN, [
            "USSD browse over the seeded curriculum",
            "SMS teaching packs, split to 160 characters",
            "Q/T question capture from any phone",
            "Topic-difficulty analytics",
            "Grade 10 Core Mathematics: 14 sub-strands",
        ]),
        ("NEXT", MUTED, [
            "English: all 45 sub-strands extracted from the\nKICD design, pages pending review",
            "The remaining compulsory learning areas",
            "Claude-generated tests from real questions",
            "Pilot with a school, measured on packs\ndelivered and questions returned",
        ]),
    ]
    for i, (label, color, rows) in enumerate(columns):
        left = MARGIN + i * Inches(6.1)
        frame = textbox(slide, left, Inches(2.35), Inches(5.6), Inches(0.35))
        write(frame, label, 13, color, bold=True, font=HEAD, first=True)
        rule(slide, left, Inches(2.75), Inches(5.6))
        frame = textbox(slide, left, Inches(2.95), Inches(5.6), Inches(2.6))
        for j, row in enumerate(rows):
            mark = "✓  " if i == 0 else "·  "
            for k, line in enumerate(row.split("\n")):
                write(frame, (mark if k == 0 else "    ") + line, 15,
                      INK if i == 0 else BODY, font=TEXT, space_after=(9 if k else 2),
                      line=1.25, first=(j == 0 and k == 0))

    band = card(slide, MARGIN, Inches(5.75), W - 2 * MARGIN, Inches(1.05), GREEN,
                GREEN)
    frame = band.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = frame.margin_right = Inches(0.35)
    write(frame, "Every Kenyan teacher already owns the hardware. We are asking for "
                 "the pilot that proves the loop.", 18, PAPER, bold=True, font=TEXT,
          align=PP_ALIGN.CENTER, first=True)
    return slide


def main():
    prs = deck()
    title_slide(prs)
    problem_slide(prs)
    how_slide(prs)
    demo_slide(prs)
    architecture_slide(prs)
    channels_slide(prs)
    ask_slide(prs)
    prs.save(OUT)
    print(f"{len(prs.slides.__iter__.__self__._sldIdLst)} slides -> {OUT}")


if __name__ == "__main__":
    main()
